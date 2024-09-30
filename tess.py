import os
import json
from PIL import Image
import pytesseract
from pdf2image import convert_from_path
from pptx import Presentation

# Specify the path to the data folder
data_folder = 'data'
output_file = 'output.jsonl'

# Function to process an image file and extract text
def process_image(image_path):
    try:
        img = Image.open(image_path)
        text = pytesseract.image_to_string(img)
        return text.strip()
    except Exception as e:
        print(f"Error processing image {image_path}: {e}")
        return None

# Function to process a PDF file and extract text from each page
def process_pdf(pdf_path):
    try:
        pages = convert_from_path(pdf_path)
        text = ''
        for page_num, page in enumerate(pages):
            print(f"Processing page {page_num + 1} of {pdf_path}")
            text += pytesseract.image_to_string(page) + '\n'
        return text.strip()
    except Exception as e:
        print(f"Error processing PDF {pdf_path}: {e}")
        return None

# Function to process a PowerPoint file and extract text
def process_pptx(pptx_path):
    try:
        presentation = Presentation(pptx_path)
        text = ''
        for slide_num, slide in enumerate(presentation.slides):
            print(f"Processing slide {slide_num + 1} of {pptx_path}")
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + '\n'
        return text.strip()
    except Exception as e:
        print(f"Error processing PowerPoint {pptx_path}: {e}")
        return None

# Function to scan all documents and write them into a JSONL file
def scan_documents_and_convert_to_jsonl(folder_path, output_file):
    with open(output_file, 'w', encoding='utf-8') as f_out:
        for filename in os.listdir(folder_path):
            file_path = os.path.join(folder_path, filename)
            print(f"Processing: {file_path}")
            
            if filename.lower().endswith(('.png', '.jpg', '.jpeg', '.tiff', '.bmp', '.gif')):
                text = process_image(file_path)
            elif filename.lower().endswith('.pdf'):
                text = process_pdf(file_path)
            elif filename.lower().endswith('.pptx'):
                text = process_pptx(file_path)
            else:
                print(f"Unsupported file format: {file_path}")
                continue  # Skip unsupported file formats
            
            if text:
                json_line = {"context": text}
                f_out.write(json.dumps(json_line) + '\n')
                #print(f"Written to JSONL: {json_line}")
            else:
                print(f"Skipped: {file_path}")

if __name__ == '__main__':
    scan_documents_and_convert_to_jsonl(data_folder, output_file)
    print(f"All documents processed and saved to {output_file}")
